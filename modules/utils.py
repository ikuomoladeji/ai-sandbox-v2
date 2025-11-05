"""
Production-grade utility functions with error handling and logging
"""
import json
import csv
import datetime
from pathlib import Path
from typing import Dict, List, Optional, Any
from reportlab.lib.pagesizes import A4
from reportlab.pdfgen import canvas
from pptx import Presentation
from pptx.util import Pt
from openpyxl import Workbook
from docx import Document

from config import Config
from modules.logger import get_logger, TPRMLogger
from modules.validators import InputValidator, ValidationError

logger = get_logger(__name__)

# Use config for paths
VENDOR_DB_PATH = Config.VENDOR_DB_PATH
HISTORY_DIR = Config.HISTORY_DIR
HISTORY_DIR.mkdir(exist_ok=True)

# ---------------- Timestamp helper ----------------

def now_iso():
    return datetime.datetime.now().strftime("%Y-%m-%d_%H-%M-%S")

def today_short():
    return datetime.datetime.now().strftime("%Y-%m-%d")

# ---------------- Risk bucket logic (centralised) ----------------

def classify_risk_bucket(likelihood, impact):
    """
    likelihood and impact are 'low' / 'medium' / 'high'
    We'll map to numbers and multiply.
    """
    mapping = {"low": 1, "medium": 2, "high": 3}
    like_score = mapping.get(likelihood, 2)
    imp_score = mapping.get(impact, 2)
    product = like_score * imp_score
    if product >= 7:
        return "high"
    elif product >= 4:
        return "medium"
    else:
        return "low"

# ---------------- Vendor DB helpers (multi-org aware) ----------------

def load_vendor_db() -> Dict[str, Any]:
    """
    Load vendor database with error handling and backup

    Returns:
        Vendor database dictionary

    Raises:
        RuntimeError: If database is corrupted and cannot be loaded
    """
    if not VENDOR_DB_PATH.exists():
        logger.info("Vendor database not found, creating new one")
        return {}

    try:
        content = VENDOR_DB_PATH.read_text(encoding="utf-8")
        db = json.loads(content)
        logger.info(f"Loaded vendor database with {len(db)} organizations")
        return db

    except json.JSONDecodeError as e:
        logger.error(f"Vendor database corrupted: {e}")

        # Try to load backup
        backup_path = VENDOR_DB_PATH.with_suffix('.json.backup')
        if backup_path.exists():
            try:
                logger.info("Attempting to load backup database")
                content = backup_path.read_text(encoding="utf-8")
                db = json.loads(content)
                logger.info("Successfully loaded backup database")
                return db
            except Exception as backup_error:
                logger.error(f"Backup database also corrupted: {backup_error}")

        raise RuntimeError(
            "Vendor database is corrupted and no valid backup found. "
            f"Please check {VENDOR_DB_PATH}"
        )

    except Exception as e:
        logger.error(f"Failed to load vendor database: {e}")
        raise RuntimeError(f"Failed to load vendor database: {e}")

def save_vendor_db(db: Dict[str, Any]) -> None:
    """
    Save vendor database with backup and validation

    Args:
        db: Vendor database to save

    Raises:
        RuntimeError: If save operation fails
    """
    try:
        # Validate database structure
        if not isinstance(db, dict):
            raise ValueError("Database must be a dictionary")

        # Create backup of existing database
        if VENDOR_DB_PATH.exists():
            backup_path = VENDOR_DB_PATH.with_suffix('.json.backup')
            try:
                VENDOR_DB_PATH.replace(backup_path)
                logger.debug("Created database backup")
            except Exception as e:
                logger.warning(f"Failed to create backup: {e}")

        # Save database
        Config.DATA_DIR.mkdir(exist_ok=True)
        content = json.dumps(db, indent=2, ensure_ascii=False)
        VENDOR_DB_PATH.write_text(content, encoding="utf-8")

        logger.info(f"Saved vendor database with {len(db)} organizations")
        TPRMLogger.log_user_action(logger, "save_vendor_db", {"org_count": len(db)})

    except Exception as e:
        logger.error(f"Failed to save vendor database: {e}")
        raise RuntimeError(f"Failed to save vendor database: {e}")

def ensure_org(db, org_id):
    if org_id not in db:
        db[org_id] = {}
    return db

def list_orgs(db):
    return list(db.keys())

def list_vendors_for_org(db, org_id):
    if org_id == "ALL":
        # flatten across all orgs
        out = []
        for org, vendors in db.items():
            for v in vendors:
                out.append((org, v))
        return out
    return [(org_id, v) for v in db.get(org_id, {}).keys()]

def snapshot_history(org_id: str, vendor_name: str, record_dict: Dict[str, Any]) -> None:
    """
    Write a point-in-time snapshot to /history with timestamp and validation

    Args:
        org_id: Organization ID
        vendor_name: Vendor name
        record_dict: Vendor record to snapshot

    Raises:
        RuntimeError: If snapshot fails
    """
    try:
        # Validate inputs
        org_id = InputValidator.validate_organization_name(org_id)
        vendor_name = InputValidator.validate_vendor_name(vendor_name)

        ts = now_iso()
        safe_org = InputValidator.sanitize_filename(org_id.replace(" ", "_"))
        safe_vendor = InputValidator.sanitize_filename(vendor_name.replace(" ", "_"))

        HISTORY_DIR.mkdir(exist_ok=True)
        hist_path = HISTORY_DIR / f"{safe_org}_{safe_vendor}_{ts}.json"

        content = json.dumps(record_dict, indent=2, ensure_ascii=False)
        hist_path.write_text(content, encoding="utf-8")

        logger.debug(f"Created history snapshot: {hist_path.name}")

    except ValidationError as e:
        logger.error(f"Validation error in snapshot_history: {e}")
        raise RuntimeError(f"Failed to create snapshot: {e}")

    except Exception as e:
        logger.error(f"Failed to create history snapshot: {e}")
        raise RuntimeError(f"Failed to create history snapshot: {e}")

def validate_vendor_record(record):
    """
    Minimal sanity checker before generating / exporting.
    Returns list of missing fields (empty list means good).
    """
    required = ["org_id", "vendor_name", "service", "business_owner",
                "overall_control_score", "likelihood", "impact"]
    missing = [f for f in required if f not in record or record[f] in ("", None)]
    return missing

# ---------------- Output format selector with isolation ----------------

def ask_output_format():
    fmt = input(
        "Output format (word/pdf/excel/powerpoint/powerbi/python/text): "
    ).strip().lower()

    if "word" in fmt or "doc" in fmt:
        return "word"
    if "pdf" in fmt:
        return "pdf"
    if "excel" in fmt or "xls" in fmt:
        return "excel"
    if "power" in fmt or "ppt" in fmt or "slide" in fmt:
        return "ppt"
    if "bi" in fmt or "powerbi" in fmt or "dashboard" in fmt:
        return "powerbi"
    if "py" in fmt or "python" in fmt:
        return "python"
    return "text"

def ask_audience_and_redaction():
    """
    We’ll tune tone + disclosure for different audiences.
    """
    audience = input(
        "Audience (internal / client / exec / vendor): "
    ).strip().lower()

    redaction = input(
        "Redaction level (full / client_safe / vendor_safe): "
    ).strip().lower()

    return audience, redaction

# ---------------- Portfolio PPT builder ----------------

def create_portfolio_ppt(org_id, db_for_report, summary_dict, outfile):
    prs = Presentation()

    # Slide 1 - Overview
    s1 = prs.slides.add_slide(prs.slide_layouts[0])
    s1.shapes.title.text = f"Third-Party Risk Overview – {org_id}"
    s1.placeholders[1].text = (
        f"Org: {org_id}\n"
        f"Total vendors: {summary_dict['total']}\n"
        f"Avg control score: {summary_dict['avg_control']}\n"
        f"High-risk: {summary_dict['risk_counts']['high']}\n"
        f"Medium-risk: {summary_dict['risk_counts']['medium']}\n"
        f"Low-risk: {summary_dict['risk_counts']['low']}\n"
        f"Weakest control area: {summary_dict['weakest_domain']}"
    )

    # Slide 2 - High Risk Vendors
    s2 = prs.slides.add_slide(prs.slide_layouts[1])
    s2.shapes.title.text = "High-Risk Vendors"
    body2 = s2.placeholders[1].text_frame
    body2.text = ""
    any_high = False
    for vendor_name, data in db_for_report.items():
        bucket = classify_risk_bucket(data["likelihood"], data["impact"])
        if bucket == "high":
            any_high = True
            p = body2.add_paragraph()
            p.text = (
                f"{vendor_name} | Score {data['overall_control_score']}/5 | "
                f"{data['likelihood']}/{data['impact']}"
            )
            p.font.size = Pt(14)
    if not any_high:
        body2.text = "No vendors currently classified as high risk."

    # Slide 3 - Common Weaknesses / Next Steps
    s3 = prs.slides.add_slide(prs.slide_layouts[1])
    s3.shapes.title.text = "Common Weaknesses / Next Steps"
    body3 = s3.placeholders[1].text_frame
    body3.text = (
        "Recurring gaps & asks:\n"
        "• Monitoring / Logging maturity\n"
        "• DR/BCP testing depth\n"
        "• Evidence freshness (SOC2, pentest)\n\n"
        "Stakeholder questions:\n"
        "• Do we have exit strategy for critical vendors?\n"
        "• Has Legal locked data residency / breach notification clauses?\n"
        "• Has the business owner accepted residual risk formally?"
    )

    prs.save(outfile)
    print(f"📊 PowerPoint deck created: {Path(outfile).resolve()}")

# ---------------- Individual vendor export ----------------

def export_vendor_document(org_id, vendor_record, narrative_text, fmt):
    """
    Exports assessment/report for a single vendor in multiple formats.
    vendor_record is the dict for this (org_id, vendor).
    """

    safe_org = vendor_record["org_id"].replace(" ", "_")
    safe_vendor = vendor_record["vendor_name"].replace(" ", "_")
    ts = today_short()
    out_dir = Path("outputs")
    out_dir.mkdir(exist_ok=True)
    base = out_dir / f"{safe_org}_{safe_vendor}_{ts}_assessment"

    # Always save baseline .txt (audit trail / reference)
    base.with_suffix(".txt").write_text(narrative_text, encoding="utf-8")

    # WORD / DOCX
    if fmt == "word":
        doc = Document()
        doc.add_heading(
            f"{vendor_record['vendor_name']} Vendor Risk Assessment ({vendor_record['org_id']})",
            0
        )
        doc.add_paragraph(narrative_text)
        doc.save(base.with_suffix(".docx"))
        print("✅ Word (.docx) created.")

    # PDF
    elif fmt == "pdf":
        style = input("PDF style (1=plain text, 2=structured sections): ").strip()
        pdf_path = base.with_suffix(".pdf")
        c = canvas.Canvas(str(pdf_path), pagesize=A4)
        textobj = c.beginText(40, 800)
        textobj.setFont("Helvetica", 10)
        for line in narrative_text.splitlines():
            textobj.textLine(line)
            if textobj.getY() < 60:
                c.drawText(textobj)
                c.showPage()
                textobj = c.beginText(40, 800)
                textobj.setFont("Helvetica", 10)
        c.drawText(textobj)
        if style == "2":
            c.showPage()
            c.drawString(40, 800, "Structured layout / next steps:")
            c.drawString(40, 780, "- Confirm DR/BCP evidence.")
            c.drawString(40, 760, "- Capture signed risk acceptance for known gaps.")
        c.save()
        print("✅ PDF created.")

    # EXCEL
    elif fmt == "excel":
        wb = Workbook()
        ws = wb.active
        ws.title = "Assessment"
        ws["A1"], ws["B1"] = "Org", vendor_record["org_id"]
        ws["A2"], ws["B2"] = "Vendor", vendor_record["vendor_name"]
        ws["A3"], ws["B3"] = "Owner", vendor_record["business_owner"]
        ws["A4"], ws["B4"] = "Likelihood", vendor_record["likelihood"]
        ws["A5"], ws["B5"] = "Impact", vendor_record["impact"]
        ws["A6"], ws["B6"] = "Overall Control Score", vendor_record["overall_control_score"]
        ws.append(["Domain", "Score"])
        for d,s in vendor_record["control_scores"].items():
            ws.append([d,s])
        # Add open actions if present
        ws2 = wb.create_sheet("Open Actions")
        ws2.append(["Owner Type","Action","Urgency","Status"])
        for a in vendor_record.get("open_actions", []):
            ws2.append([
                a.get("owner_type",""),
                a.get("action",""),
                a.get("urgency",""),
                a.get("status","open"),
            ])
        wb.save(base.with_suffix(".xlsx"))
        print("✅ Excel (.xlsx) created.")

    # POWERPOINT
    elif fmt == "ppt":
        kind = input("PowerPoint type (1=Exec summary, 2=Vendor deep dive): ").strip()
        prs = Presentation()

        # Slide 1: summary
        s1 = prs.slides.add_slide(prs.slide_layouts[0])
        s1.shapes.title.text = f"Vendor Risk Overview – {vendor_record['vendor_name']} ({vendor_record['org_id']})"
        s1.placeholders[1].text = (
            f"Overall control score: {vendor_record['overall_control_score']}/5\n"
            f"Likelihood: {vendor_record['likelihood']}\n"
            f"Impact: {vendor_record['impact']}\n"
            f"Risk Bucket: {classify_risk_bucket(vendor_record['likelihood'], vendor_record['impact'])}\n"
        )

        if kind == "2":
            # Slide 2: Control domains
            s2 = prs.slides.add_slide(prs.slide_layouts[1])
            s2.shapes.title.text = "Control Domains"
            body2 = s2.placeholders[1].text_frame
            body2.text = "Control strength by area:"
            for d,s in vendor_record["control_scores"].items():
                p = body2.add_paragraph()
                p.text = f"{d}: {s}/5"
                p.font.size = Pt(14)

            # Slide 3: Recommendations / Actions
            s3 = prs.slides.add_slide(prs.slide_layouts[1])
            s3.shapes.title.text = "Open Actions / Recommendations"
            body3 = s3.placeholders[1].text_frame
            body3.text = ""
            for a in vendor_record.get("open_actions", []):
                p = body3.add_paragraph()
                p.text = f"[{a.get('urgency','?').upper()}] {a.get('owner_type','')}: {a.get('action','')}"
                p.font.size = Pt(14)

            # Slide 4: Narrative (trimmed)
            s4 = prs.slides.add_slide(prs.slide_layouts[1])
            s4.shapes.title.text = "Narrative Summary"
            body4 = s4.placeholders[1].text_frame
            body4.text = narrative_text[:2000]

        prs.save(base.with_suffix(".pptx"))
        print("✅ PowerPoint (.pptx) created.")

    # POWER BI FEED
    elif fmt == "powerbi":
        pbi_dir = Path("outputs/powerbi")
        pbi_dir.mkdir(exist_ok=True)
        csv_path = pbi_dir / f"{safe_org}_{safe_vendor}_{ts}.csv"
        with open(csv_path, "w", newline="", encoding="utf-8") as f:
            writer = csv.writer(f)
            writer.writerow(["Org","Vendor","Domain","Score","Overall","Likelihood","Impact","RiskBucket"])
            bucket = classify_risk_bucket(vendor_record["likelihood"], vendor_record["impact"])
            for d,s in vendor_record["control_scores"].items():
                writer.writerow([
                    vendor_record["org_id"],
                    vendor_record["vendor_name"],
                    d,
                    s,
                    vendor_record["overall_control_score"],
                    vendor_record["likelihood"],
                    vendor_record["impact"],
                    bucket
                ])
        print("✅ PowerBI CSV created.")

    # PYTHON
    elif fmt == "python":
        py_path = base.with_suffix(".py")
        py_code = "vendor_data = " + json.dumps(vendor_record, indent=2)
        py_path.write_text(py_code, encoding="utf-8")
        print("✅ Python .py data created.")

    # TEXT or fallback
    else:
        print("Saved baseline .txt only.")

# ---------------- Portfolio export helpers ----------------

def export_portfolio_artifacts(org_id, db_for_report, summary_dict, narrative_text, fmt):
    """
    org_id can be a real org name or 'ALL'.
    db_for_report is { vendor_name: vendor_record, ... } for that org.
    """

    ts = today_short()
    safe_org = org_id.replace(" ", "_")
    out_dir = Path("outputs")
    out_dir.mkdir(exist_ok=True)
    base = out_dir / f"Portfolio_{safe_org}_{ts}"

    # always write the narrative text
    base.with_suffix(".txt").write_text(narrative_text, encoding="utf-8")

    if fmt == "powerbi":
        powerbi_dir = Path("outputs/powerbi")
        powerbi_dir.mkdir(exist_ok=True)
        csv_path = powerbi_dir / f"portfolio_{safe_org}_{ts}.csv"
        with open(csv_path,"w",newline="",encoding="utf-8") as f:
            writer = csv.writer(f)
            writer.writerow([
                "Org","Vendor","OverallControlScore","Likelihood","Impact","RiskBucket"
            ])
            for vendor_name, v in db_for_report.items():
                bucket = classify_risk_bucket(v["likelihood"], v["impact"])
                writer.writerow([
                    v["org_id"],
                    v["vendor_name"],
                    v["overall_control_score"],
                    v["likelihood"],
                    v["impact"],
                    bucket
                ])
        print(f"Saved PowerBI dataset: {csv_path.resolve()}")

    elif fmt == "ppt":
        ppt_path = base.with_suffix(".pptx")
        create_portfolio_ppt(org_id, db_for_report, summary_dict, ppt_path)

    elif fmt == "python":
        py_path = base.with_suffix(".py")
        py_payload = {
            "summary": summary_dict,
            "vendors": db_for_report
        }
        py_path.write_text("portfolio_summary = " + json.dumps(py_payload, indent=2),
                           encoding="utf-8")
        print(f"Saved Python summary module: {py_path.resolve()}")

    elif fmt == "excel":
        wb = Workbook()
        ws = wb.active
        ws.title = "Portfolio Summary"
        ws["A1"] = f"Org: {org_id}"
        ws["A2"] = f"Total vendors: {summary_dict['total']}"
        ws["A3"] = f"Avg control score: {summary_dict['avg_control']}"
        ws["A4"] = f"High risk: {summary_dict['risk_counts']['high']}"
        ws["A5"] = f"Weakest domain: {summary_dict['weakest_domain']}"
        ws2 = wb.create_sheet("Vendor Risk")
        ws2.append(["Vendor","Owner","OverallScore","Likelihood","Impact","RiskBucket"])
        for vendor_name, v in db_for_report.items():
            bucket = classify_risk_bucket(v["likelihood"], v["impact"])
            ws2.append([
                vendor_name,
                v["business_owner"],
                v["overall_control_score"],
                v["likelihood"],
                v["impact"],
                bucket
            ])
        xlsx_path = base.with_suffix(".xlsx")
        wb.save(xlsx_path)
        print(f"Saved Excel workbook: {xlsx_path.resolve()}")

    elif fmt == "word":
        doc = Document()
        doc.add_heading(f"Portfolio / Management Report – {org_id}", 0)
        doc.add_paragraph(narrative_text)
        doc.save(base.with_suffix(".docx"))
        print(f"Saved Word report: {base.with_suffix('.docx').resolve()}")

    elif fmt == "pdf":
        pdf_path = base.with_suffix(".pdf")
        c = canvas.Canvas(str(pdf_path), pagesize=A4)
        textobj = c.beginText(40, 800)
        textobj.setFont("Helvetica", 10)
        for line in narrative_text.splitlines():
            textobj.textLine(line)
            if textobj.getY() < 60:
                c.drawText(textobj)
                c.showPage()
                textobj = c.beginText(40, 800)
                textobj.setFont("Helvetica", 10)
        c.drawText(textobj)
        c.save()
        print(f"Saved PDF report: {pdf_path.resolve()}")

    else:
        print("Saved baseline .txt narrative only.")
